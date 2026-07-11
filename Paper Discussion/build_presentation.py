"""
Build the Arabic Text Correction paper presentation from the IMSA template.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
TEMPLATE = ROOT / "Paper_Presentation_Template.pptx"
OUT = ROOT / "Arabic_Text_Correction_Presentation.pptx"
ASSETS = ROOT / "assets"
DIAGRAMS = ROOT.parent / "paper" / "Diagrams"
TEMPLATE_ASSETS = ROOT / "_template_assets"

BLUE = RGBColor(0x00, 0x64, 0x99)
BLUE_DARK = RGBColor(0x00, 0x44, 0x66)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x4A, 0x5A, 0x6A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF0, 0xF7, 0xFB)
ACCENT = RGBColor(0xE8, 0xA8, 0x38)


def delete_all_slides(prs: Presentation) -> None:
    """Remove every slide while keeping slide master / theme."""
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(qn("r:id"))
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def set_run(run, text, size=18, bold=False, color=BLACK, font="Arial"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = color


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT, font="Arial"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    # Support multi-line text via explicit newlines
    lines = text.split("\n")
    for i, line in enumerate(lines):
        para = p if i == 0 else tf.add_paragraph()
        para.alignment = align
        run = para.add_run()
        set_run(run, line, size=size, bold=bold, color=color, font=font)
    return box


def fill_title(slide, title: str, size=32):
    """Fill the title placeholder if present, else add a textbox."""
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx == 0:
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            set_run(run, title, size=size, bold=True, color=BLUE, font="Arial")
            return shape
    return add_textbox(slide, Inches(0.66), Inches(0.28), Inches(11.5), Inches(0.9), title, size=size, bold=True, color=BLUE)


def add_bullets(slide, left, top, width, height, items, size=18, color=BLACK, bold_first=False, spacing=8):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        p.space_after = Pt(spacing)
        # bullet via text prefix for reliable rendering
        text = item if item.startswith("•") or item.startswith("–") else f"• {item}"
        run = p.add_run()
        is_bold = bold_first and i == 0
        set_run(run, text, size=size, bold=is_bold, color=color, font="Arial")
    return box


def add_picture_safe(slide, path, left, top, width=None, height=None):
    path = Path(path)
    if not path.exists():
        print("MISSING IMAGE:", path)
        return None
    kwargs = {}
    if width is not None:
        kwargs["width"] = width
    if height is not None:
        kwargs["height"] = height
    return slide.shapes.add_picture(str(path), left, top, **kwargs)


def add_logo_corner(slide, prs):
    """Add IMSA logo top-right like template content slides."""
    logo = TEMPLATE_ASSETS / "Picture 2.png"
    if logo.exists():
        add_picture_safe(slide, logo, Inches(9.05), Inches(0.05), width=Inches(3.5))


def add_footer_note(slide, text, left=Inches(0.5), top=Inches(7.05), width=Inches(12)):
    add_textbox(slide, left, top, width, Inches(0.35), text, size=10, color=GRAY, font="Arial")


def add_card(slide, left, top, width, height, fill=LIGHT_BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = BLUE
    shape.line.width = Pt(1.25)
    return shape


def add_notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


# ---------- slide builders ----------

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # conference header
    add_textbox(
        slide,
        Inches(0.4),
        Inches(0.25),
        Inches(8.5),
        Inches(0.8),
        "Intelligent Methods, Systems, and Applications\n(IMSA)",
        size=16,
        bold=True,
        color=BLUE,
        font="Noto Serif",
    )
    logo = TEMPLATE_ASSETS / "Picture 2.png"
    if logo.exists():
        add_picture_safe(slide, logo, Inches(9.0), Inches(0.0), width=Inches(3.6))

    add_textbox(
        slide,
        Inches(0.6),
        Inches(1.4),
        Inches(11.5),
        Inches(1.4),
        "An End-to-End Arabic Text Correction Framework:\nComparing Custom Character-Level Transformers\nand Pre-trained AraBART",
        size=26,
        bold=True,
        color=BLUE,
        align=PP_ALIGN.CENTER,
    )

    authors = [
        ("Mohamed Soltan", "Dept. of Computer Science\nMSA University"),
        ("Abdelmonem Hatem", "Dept. of Computer Science\nMSA University"),
        ("Mohamed Taha", "Dept. of Computer Science\nMSA University"),
        ("Youssef Khalaf", "Dept. of Computer Science\nMSA University"),
    ]
    xs = [0.4, 3.5, 6.6, 9.7]
    for x, (name, aff) in zip(xs, authors):
        add_textbox(slide, Inches(x), Inches(3.3), Inches(2.9), Inches(0.4), name, size=13, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x), Inches(3.7), Inches(2.9), Inches(0.7), aff, size=11, color=GRAY, align=PP_ALIGN.CENTER)

    msa = TEMPLATE_ASSETS / "Picture 6.jpg"
    if msa.exists():
        add_picture_safe(slide, msa, Inches(0.35), Inches(5.9), width=Inches(1.1))

    add_textbox(
        slide,
        Inches(1.7),
        Inches(6.2),
        Inches(7),
        Inches(0.6),
        "October University for Modern Sciences and Arts (MSA)\nCairo, Egypt",
        size=12,
        color=GRAY,
    )
    add_textbox(slide, Inches(9.5), Inches(6.5), Inches(2.8), Inches(0.4), "IEEE / IMSA 2026", size=14, bold=True, color=BLUE, align=PP_ALIGN.RIGHT)
    add_notes(
        slide,
        "Good morning. We present an end-to-end Arabic text correction framework comparing a compact custom Transformer with a pre-trained AraBART multi-modal service.",
    )
    return slide


def slide_outline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_title(slide, "Outline")
    add_logo_corner(slide, prs)
    items = [
        "Introduction & Motivation",
        "Problem Statement & Soft Errors",
        "Related Work",
        "Proposed End-to-End Framework",
        "Data Pipeline, Normalization & Noise Injection",
        "Custom Transformer Architecture",
        "AraBART Multi-modal Deployment Service",
        "Experiments, Metrics & Baselines",
        "Results, Error Analysis & Discussion",
        "Conclusion & Future Work",
    ]
    add_bullets(slide, Inches(1.0), Inches(1.55), Inches(10), Inches(5.4), items, size=20, spacing=10)
    notes = slide.notes_slide.notes_text_frame
    notes.text = "Walk through the agenda briefly. Emphasize the two-path story: controlled benchmarking vs. real-world deployment."
    return slide


def slide_intro(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_title(slide, "Introduction")
    add_logo_corner(slide, prs)
    add_bullets(
        slide,
        Inches(0.6),
        Inches(1.4),
        Inches(6.2),
        Inches(5.2),
        [
            "Arabic text correction is critical for education, media QA, accessibility, and moderation.",
            "Challenges: orthographic variants, visually confusable letters, optional diacritics, and spacing that changes meaning.",
            "OCR and ASR add broken words, missing characters, merged tokens, and repeated substitutions.",
            "Core design question: compact task-specific model vs. larger pre-trained Arabic model?",
            "This work studies that trade-off with two complementary correction paths.",
        ],
        size=16,
        spacing=12,
    )
    add_picture_safe(slide, ASSETS / "soft_errors.png", Inches(6.8), Inches(1.5), width=Inches(5.7))
    return slide


def slide_motivation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_title(slide, "Motivation")
    add_logo_corner(slide, prs)
    cards = [
        ("OCR Noise", "Broken tokens, substitutions, and spacing errors from scanned Arabic documents."),
        ("ASR Noise", "Speech transcripts introduce phonetic and boundary mistakes."),
        ("Typing Soft Errors", "Alef/Hamza/Teh Marbuta confusions are frequent and meaning-changing."),
        ("Deployment Gap", "Most Arabic studies report scores, few ship multi-modal user-facing systems."),
    ]
    positions = [(0.5, 1.5), (6.5, 1.5), (0.5, 4.0), (6.5, 4.0)]
    for (x, y), (title, body) in zip(positions, cards):
        add_card(slide, Inches(x), Inches(y), Inches(5.7), Inches(2.1))
        add_textbox(slide, Inches(x + 0.25), Inches(y + 0.25), Inches(5.2), Inches(0.45), title, size=18, bold=True, color=BLUE)
        add_textbox(slide, Inches(x + 0.25), Inches(y + 0.85), Inches(5.2), Inches(1.0), body, size=14, color=BLACK)
    return slide


def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_title(slide, "Problem Statement")
    add_logo_corner(slide, prs)
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.5),
        Inches(11.5),
        Inches(5.0),
        [
            "Small character-level errors in Arabic can change meaning, break token boundaries, or reduce grammatical fluency.",
            "Correction is more than dictionary lookup — it needs sequence context at character and sentence levels.",
            "Existing resources often lack controlled noisy-clean pairs tailored to modern Arabic news orthography.",
            "Benchmarks alone are insufficient: real systems must also handle OCR images, audio, and live speech.",
            "Goal: a reproducible end-to-end workflow covering data curation, noise synthesis, training, evaluation, and deployment.",
        ],
        size=17,
        spacing=14,
    )
    return slide


def slide_soft_errors(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_title(slide, "Arabic Soft Error Categories")
    add_logo_corner(slide, prs)
    add_picture_safe(slide, ASSETS / "soft_errors.png", Inches(0.6), Inches(1.35), width=Inches(12.0))
    add_footer_note(slide, "Guided by learner studies, social-media patterns, and editor correction logs (~85% of common Arabic spelling errors).")
    return slide


def slide_contributions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_title(slide, "Contributions")
    add_logo_corner(slide, prs)
    add_picture_safe(slide, ASSETS / "contributions.png", Inches(0.7), Inches(1.5), width=Inches(11.8))
    return slide


def slide_related(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_title(slide, "Related Work")
    add_logo_corner(slide, prs)
    rows = [
        ("Seq2Seq / Transformers", "Sutskever, Bahdanau, Vaswani — correction as conditional generation with attention."),
        ("Denoising Pretraining", "BART / T5 show corruption-reconstruction transfers well to correction tasks."),
        ("Arabic Pretraining", "AraBERT & CAMeL models improve Arabic understanding and generation baselines."),
        ("Arabic GEC Benchmarks", "QALB remains core; BiLSTM and T5 variants still struggle with heavy noise."),
        ("Deployment Gap", "Few works combine OCR + ASR + correction in one Arabic user-facing application."),
    ]
    y = 1.4
    for title, body in rows:
        add_textbox(slide, Inches(0.7), Inches(y), Inches(3.6), Inches(0.7), title, size=14, bold=True, color=BLUE)
        add_textbox(slide, Inches(4.4), Inches(y), Inches(8.0), Inches(0.7), body, size=14, color=BLACK)
        y += 0.95
    return slide


def slide_framework(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_title(slide, "Proposed Framework Overview")
    add_logo_corner(slide, prs)
    add_picture_safe(slide, ASSETS / "two_tier.png", Inches(0.5), Inches(1.25), width=Inches(12.2))
    add_notes(
        slide,
        "This is the key conceptual slide: Branch A optimizes controlled character fidelity; Branch B optimizes real-world multi-modal usability.",
    )
    return slide


def slide_pipeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_title(slide, "End-to-End Pipeline")
    add_logo_corner(slide, prs)
    img = DIAGRAMS / "pipeline.png"
    add_picture_safe(slide, img, Inches(0.25), Inches(1.2), width=Inches(12.7))
    add_footer_note(slide, "Shared data path splits into Branch A (custom Transformer) and Branch B (AraBART multi-modal service).")
    add_notes(slide, "Walk left-to-right: Youm7 corpus, cleaning, normalization, synthetic noise, then the two branches.")
    return slide


def slide_data(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_title(slide, "Data Acquisition & Curation")
    add_logo_corner(slide, prs)
    add_bullets(
        slide,
        Inches(0.6),
        Inches(1.4),
        Inches(6.3),
        Inches(5.0),
        [
            "Scraped Youm7 Arabic news with async aiohttp + BeautifulSoup.",
            "Polite crawling: 1–3s randomized delays, retries, robots.txt compliance.",
            "Filters: remove boilerplate, short articles, and samples with >15% non-Arabic.",
            "Final clean corpus: 100,000 high-quality articles (news, education, health, tech).",
            "Prepared 10,000 noisy-clean pairs for controlled modeling.",
            "Split: 8,000 train / 1,000 val / 1,000 held-out test.",
        ],
        size=15,
        spacing=10,
    )
    add_picture_safe(slide, ASSETS / "dataset_split.png", Inches(7.0), Inches(1.5), width=Inches(5.4))
    return slide


def slide_corpus_table(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_title(slide, "Arabic Corpus Context")
    add_logo_corner(slide, prs)

    rows = [
        ["Dataset", "Scale", "Source", "Use / Availability"],
        ["Youm7 (this work)", "100,000", "Youm7 Egypt", "Noisy-clean generation / Private"],
        ["SANAD", "194,797", "AlKhaleej+", "Classification / Public"],
        ["Ultimate Arabic News", "193,000", "Multiple outlets", "General NLP / Public"],
        ["ANAD", "500,000", "12 news sites", "Annotated news / Public"],
        ["Amina", "~1.85M", "Regional newspapers", "Multimodal metadata / Public"],
        ["Wiki-40B (AR)", "~1M+", "Arabic Wikipedia", "Language modeling / Public"],
    ]
    table = slide.shapes.add_table(len(rows), 4, Inches(0.45), Inches(1.45), Inches(12.3), Inches(5.2)).table
    widths = [Inches(3.0), Inches(1.8), Inches(2.8), Inches(4.7)]
    for i, w in enumerate(widths):
        table.columns[i].width = w
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(12 if r else 13)
                    run.font.bold = r == 0 or (r == 1 and c == 0)
                    run.font.name = "Arial"
                    run.font.color.rgb = WHITE if r == 0 else BLACK
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BLUE
            elif r == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG
    return slide


def slide_noise(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_title(slide, "Normalization & Noise Injection")
    add_logo_corner(slide, prs)
    add_bullets(
        slide,
        Inches(0.5),
        Inches(1.35),
        Inches(6.0),
        Inches(3.2),
        [
            "Deterministic normalization: Alef, Yaa, Hamza, Teh Marbuta; remove tatweel/diacritics where needed.",
            "Overall corruption budget: 20% (10% sub, 5% del, 5% ins).",
            "Confusable-character map + keyboard adjacency + punctuation drift.",
            "Categories keep synthesis realistic and analyzable.",
        ],
        size=14,
        spacing=8,
    )
    add_picture_safe(slide, ASSETS / "corruption_budget.png", Inches(6.5), Inches(1.3), width=Inches(6.0))

    # example pairs
    add_textbox(slide, Inches(0.5), Inches(4.7), Inches(12), Inches(0.35), "Representative Generated Pairs", size=14, bold=True, color=BLUE)
    add_textbox(
        slide,
        Inches(0.5),
        Inches(5.1),
        Inches(12.2),
        Inches(1.6),
        "Noisy: اعلنت كليه الصيدله عن مواعبد التسجيل\nClean: أعلنت كلية الصيدلة عن مواعيد التسجيل\n\nNoisy: انلقت مبادرة مدارس النيل المصرية الودلية\nClean: انطلقت مبادرة مدارس النيل المصرية الدولية",
        size=13,
        color=BLACK,
    )
    return slide


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_title(slide, "Custom Seq2Seq Transformer (Branch A)")
    add_logo_corner(slide, prs)
    add_picture_safe(slide, DIAGRAMS / "seq2seq.png", Inches(0.2), Inches(1.15), width=Inches(12.8))
    return slide


def slide_hyperparams(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_title(slide, "Model Configuration & Training")
    add_logo_corner(slide, prs)

    left_items = [
        "Character-level encoder-decoder Transformer",
        "Embedding dim: 256",
        "Encoder / Decoder layers: 3 / 3",
        "Attention heads: 8",
        "FFN dimension: 512",
        "Dropout: 0.1",
        "Max sequence length: 128",
    ]
    right_items = [
        "Optimizer: Adam (lr = 1e-4)",
        "Loss: sparse categorical cross-entropy",
        "Scheduled sampling: teacher forcing 1.0 → 0.3",
        "Mixed precision + checkpointing",
        "Batch size: 32 | Epochs: 10",
        "GPU training ≈ 100 minutes",
        "Checkpoint size ≈ 50 MB",
    ]
    add_card(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.2))
    add_textbox(slide, Inches(0.75), Inches(1.55), Inches(5.3), Inches(0.4), "Architecture", size=18, bold=True, color=BLUE)
    add_bullets(slide, Inches(0.75), Inches(2.1), Inches(5.3), Inches(4.2), left_items, size=15, spacing=8)

    add_card(slide, Inches(6.7), Inches(1.4), Inches(5.8), Inches(5.2))
    add_textbox(slide, Inches(6.95), Inches(1.55), Inches(5.3), Inches(0.4), "Training Protocol", size=18, bold=True, color=BLUE)
    add_bullets(slide, Inches(6.95), Inches(2.1), Inches(5.3), Inches(4.2), right_items, size=15, spacing=8)
    return slide


def slide_arabart(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_title(slide, "AraBART Deployment Service (Branch B)")
    add_logo_corner(slide, prs)
    add_bullets(
        slide,
        Inches(0.6),
        Inches(1.4),
        Inches(6.2),
        Inches(5.0),
        [
            "Model: CAMeL-Lab/arabart-qalb15-gec-ged-13 (Hugging Face).",
            "Integrated in a Streamlit application for practical use.",
            "Input modes: manual text, text files, OCR images, audio, live speech.",
            "OCR via OCR.space API; speech via Whisper ASR.",
            "Zero-shot deployment setting — no project-specific fine-tuning in this study.",
            "Stronger grammatical fluency and broader real-world usability.",
        ],
        size=15,
        spacing=10,
    )
    add_picture_safe(slide, ASSETS / "footprint.png", Inches(6.9), Inches(1.5), width=Inches(5.5))
    return slide


def slide_metrics(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_title(slide, "Evaluation Metrics")
    add_logo_corner(slide, prs)
    add_card(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(3.2))
    add_textbox(slide, Inches(0.9), Inches(1.7), Inches(5.2), Inches(0.4), "Primary Metric: CER", size=18, bold=True, color=BLUE)
    add_textbox(
        slide,
        Inches(0.9),
        Inches(2.3),
        Inches(5.2),
        Inches(2.0),
        "CER = (S + D + I) / N\n\nS = substitutions\nD = deletions\nI = insertions\nN = reference length (characters)",
        size=16,
        color=BLACK,
    )
    add_bullets(
        slide,
        Inches(6.8),
        Inches(1.5),
        Inches(5.5),
        Inches(4.5),
        [
            "CER suits character-level soft spelling correction.",
            "BLEU used as secondary lexical/sequence quality metric.",
            "Two evaluation settings (not a single leaderboard):",
            "– Custom model: in-domain synthetic held-out split",
            "– AraBART: zero-shot multi-modal application inputs",
            "Baselines: normalization-only + dictionary spell correction;",
            "project reference ranges for rule-based and BiLSTM.",
        ],
        size=14,
        spacing=8,
    )
    return slide


def slide_baselines(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_title(slide, "Baseline Context")
    add_logo_corner(slide, prs)
    add_picture_safe(slide, ASSETS / "baseline_bars.png", Inches(1.2), Inches(1.3), width=Inches(10.5))
    add_footer_note(slide, "Custom Transformer reaches 89.36% character accuracy on synthetic project runs — best in-domain fidelity.")
    return slide


def slide_training_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_title(slide, "Training Dynamics")
    add_logo_corner(slide, prs)
    add_picture_safe(slide, ASSETS / "training_curve.png", Inches(0.7), Inches(1.25), width=Inches(8.0))
    add_bullets(
        slide,
        Inches(8.9),
        Inches(1.6),
        Inches(3.5),
        Inches(4.8),
        [
            "Epoch 1: 38.38%",
            "Epoch 6 peak: 91.69%",
            "Epoch 10: 89.36%",
            "Fast early learning of dominant patterns.",
            "Later slight drop suggests mild overfitting to synthetic noise.",
            "Earlier stopping / stronger regularization recommended.",
        ],
        size=13,
        spacing=8,
    )
    return slide


def slide_results_table(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_title(slide, "Side-by-Side Results")
    add_logo_corner(slide, prs)

    rows = [
        ["Attribute", "Custom Transformer", "AraBART Service"],
        ["Setting", "In-domain synthetic", "Zero-shot multi-modal"],
        ["Model", "Char seq2seq Transformer", "arabart-qalb15-gec-ged-13"],
        ["CER", "0.0364", "0.0950"],
        ["BLEU", "0.292", "Not reported"],
        ["Peak token acc.", "91.69%", "N/A (pre-trained)"],
        ["Fine-tuning", "Yes", "No"],
        ["Training time", "~100 min GPU", "Inference only"],
        ["Size", "~50 MB", "~1.5 GB"],
        ["Inputs", "Text pipeline", "Text, file, OCR, audio, speech"],
    ]
    table = slide.shapes.add_table(len(rows), 3, Inches(0.5), Inches(1.3), Inches(12.2), Inches(5.5)).table
    widths = [Inches(2.6), Inches(4.8), Inches(4.8)]
    for i, w in enumerate(widths):
        table.columns[i].width = w
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(12 if r else 13)
                    run.font.bold = r == 0 or c == 0
                    run.font.name = "Arial"
                    run.font.color.rgb = WHITE if r == 0 else BLACK
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BLUE
            elif r % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG
    add_notes(
        slide,
        "Stress that these are two operating settings, not a single fair leaderboard. Custom wins CER in-domain; AraBART wins deployment breadth.",
    )
    return slide


def slide_cer_compare(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_title(slide, "CER Comparison Across Settings")
    add_logo_corner(slide, prs)
    add_picture_safe(slide, ASSETS / "cer_comparison.png", Inches(0.8), Inches(1.3), width=Inches(7.8))
    add_bullets(
        slide,
        Inches(8.8),
        Inches(1.8),
        Inches(3.6),
        Inches(4.5),
        [
            "Lower CER = better reconstruction.",
            "Custom model wins in-domain fidelity.",
            "AraBART remains preferred for fluency and multi-modal robustness.",
            "Not a like-for-like contest — two operating settings.",
        ],
        size=14,
        spacing=10,
    )
    return slide


def slide_cer_dist(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_title(slide, "CER Distribution")
    add_logo_corner(slide, prs)
    add_picture_safe(slide, DIAGRAMS / "cer.png", Inches(0.5), Inches(1.3), width=Inches(8.2))
    add_bullets(
        slide,
        Inches(8.9),
        Inches(1.6),
        Inches(3.5),
        Inches(5.0),
        [
            "Most samples concentrate in a low-CER band.",
            "Indicates stable correction, not luck on a few cases.",
            "Right tail = heavy corruption samples.",
            "Useful for reliability assessment in deployment.",
        ],
        size=14,
        spacing=10,
    )
    return slide


def slide_confusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_title(slide, "Character Confusion Matrix")
    add_logo_corner(slide, prs)
    add_picture_safe(slide, DIAGRAMS / "cm.png", Inches(0.3), Inches(1.15), height=Inches(5.8))
    add_bullets(
        slide,
        Inches(7.0),
        Inches(1.5),
        Inches(5.5),
        Inches(5.2),
        [
            "Strong diagonal = solid identity mapping.",
            "Off-diagonal cells reveal persistent substitutions.",
            "Hard cases: Alef/Hamza variants, Yaa/Alif-Maqsura, Teh Marbuta/Heh.",
            "Whitespace boundary errors amplify local drift.",
            "Guides targeted augmentation and loss design.",
        ],
        size=14,
        spacing=10,
    )
    return slide


def slide_error_analysis(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_title(slide, "Error Analysis")
    add_logo_corner(slide, prs)
    add_picture_safe(slide, ASSETS / "error_pie.png", Inches(0.4), Inches(1.3), width=Inches(6.0))
    add_bullets(
        slide,
        Inches(6.7),
        Inches(1.5),
        Inches(5.7),
        Inches(5.2),
        [
            "Substitutions dominate (~60–70%).",
            "Insertions and deletions each ~15–20%.",
            "Matches Arabic soft-error behavior: near-equivalent character choice.",
            "Future work should prioritize whitespace + Alef/Hamza contrasts.",
            "Curriculum / class-weighted objectives can reduce hard confusions.",
        ],
        size=15,
        spacing=10,
    )
    return slide


def slide_discussion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_title(slide, "Discussion & Deployment Implications")
    add_logo_corner(slide, prs)
    add_card(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.0))
    add_textbox(slide, Inches(0.75), Inches(1.6), Inches(5.3), Inches(0.4), "When to use Custom Transformer", size=16, bold=True, color=BLUE)
    add_bullets(
        slide,
        Inches(0.75),
        Inches(2.2),
        Inches(5.3),
        Inches(4.0),
        [
            "Controlled in-domain character fidelity",
            "Low footprint (~50 MB) after training",
            "Interpretable error diagnostics",
            "Edge / lightweight serving",
        ],
        size=14,
        spacing=8,
    )

    add_card(slide, Inches(6.7), Inches(1.4), Inches(5.8), Inches(5.0), fill=RGBColor(0xFF, 0xF8, 0xEC))
    add_textbox(slide, Inches(6.95), Inches(1.6), Inches(5.3), Inches(0.4), "When to use AraBART Service", size=16, bold=True, color=ACCENT)
    add_bullets(
        slide,
        Inches(6.95),
        Inches(2.2),
        Inches(5.3),
        Inches(4.0),
        [
            "Grammatical fluency under real noise",
            "OCR / ASR / live speech inputs",
            "Fast adoption (no fine-tuning here)",
            "User-facing educational & media tools",
        ],
        size=14,
        spacing=8,
    )
    return slide


def slide_conclusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_title(slide, "Conclusion")
    add_logo_corner(slide, prs)
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.5),
        Inches(11.5),
        Inches(5.2),
        [
            "Delivered a reproducible Arabic correction workflow: collect → normalize → synthesize → train → evaluate → deploy.",
            "Custom character-level Transformer: strong in-domain reconstruction (CER 0.0364, peak accuracy 91.69%).",
            "AraBART Streamlit service: better practical fluency across text, OCR, audio, and live speech.",
            "Remaining errors concentrate in whitespace and Alef/Hamza-style confusions.",
            "Best practice: two-tier system — compact model for control, pre-trained service for real-world use.",
        ],
        size=16,
        spacing=12,
    )
    add_notes(
        slide,
        "Close with the two-tier takeaway and the three contribution pillars: workflow, dual paths, and strong in-domain CER.",
    )
    return slide


def slide_future(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_title(slide, "Future Work")
    add_logo_corner(slide, prs)
    items = [
        ("Targeted Augmentation", "Focus synthetic noise on Alef/Hamza, Teh Marbuta/Heh, and whitespace boundaries."),
        ("Regularization & Early Stop", "Preserve peak generalization near epoch 6 and reduce late overfitting."),
        ("Shared Benchmarks", "Fair comparison on external datasets such as QALB."),
        ("Dialect Coverage", "Extend to dialectal Arabic with dialect ID + focused data collection."),
        ("Grammar Layer", "Add grammar-aware detection for broader writing-support systems."),
        ("Confidence Estimation", "Production-ready uncertainty signals for safer corrections."),
    ]
    positions = [(0.5, 1.4), (6.5, 1.4), (0.5, 3.2), (6.5, 3.2), (0.5, 5.0), (6.5, 5.0)]
    for (x, y), (title, body) in zip(positions, items):
        add_card(slide, Inches(x), Inches(y), Inches(5.7), Inches(1.55))
        add_textbox(slide, Inches(x + 0.2), Inches(y + 0.15), Inches(5.3), Inches(0.35), title, size=14, bold=True, color=BLUE)
        add_textbox(slide, Inches(x + 0.2), Inches(y + 0.55), Inches(5.3), Inches(0.85), body, size=12, color=BLACK)
    return slide


def slide_refs(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_title(slide, "Selected References")
    add_logo_corner(slide, prs)
    refs = [
        "Vaswani et al., Attention Is All You Need, NeurIPS, 2017.",
        "Lewis et al., BART: Denoising Seq2Seq Pre-training, ACL, 2020.",
        "Raffel et al., Exploring the Limits of Transfer Learning with T5, JMLR, 2020.",
        "Antoun et al., AraBERT: Transformer-based Model for Arabic, 2020.",
        "Abandah et al., Correcting Arabic Soft Spelling Mistakes Using BiLSTM, 2022.",
        "Al-Qaraghuli & Jaafar, Arabic Soft Spelling Correction with T5, 2024.",
        "Rozovskaya et al., QALB Arabic Error Correction Shared Tasks.",
        "CAMeL-Lab AraBART (arabart-qalb15-gec-ged-13), Hugging Face.",
        "Radford et al., Whisper: Robust Speech Recognition, 2023.",
    ]
    add_bullets(slide, Inches(0.7), Inches(1.4), Inches(11.5), Inches(5.5), refs, size=14, spacing=8)
    return slide


def slide_thanks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    logo = TEMPLATE_ASSETS / "Picture 2.png"
    if logo.exists():
        add_picture_safe(slide, logo, Inches(9.0), Inches(0.1), width=Inches(3.5))
    add_textbox(
        slide,
        Inches(1),
        Inches(2.2),
        Inches(11),
        Inches(1.2),
        "Thank You",
        size=48,
        bold=True,
        color=BLUE,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Inches(1),
        Inches(3.5),
        Inches(11),
        Inches(0.6),
        "Questions & Discussion",
        size=24,
        color=BLUE_DARK,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Inches(1),
        Inches(4.5),
        Inches(11),
        Inches(1.2),
        "Mohamed Soltan  •  Abdelmonem Hatem  •  Mohamed Taha  •  Youssef Khalaf\nOctober University for Modern Sciences and Arts (MSA)",
        size=14,
        color=GRAY,
        align=PP_ALIGN.CENTER,
    )
    msa = TEMPLATE_ASSETS / "Picture 6.jpg"
    if msa.exists():
        add_picture_safe(slide, msa, Inches(5.7), Inches(5.9), width=Inches(1.2))
    return slide


def build():
    prs = Presentation(str(TEMPLATE))
    delete_all_slides(prs)

    builders = [
        slide_title,
        slide_outline,
        slide_intro,
        slide_motivation,
        slide_problem,
        slide_soft_errors,
        slide_contributions,
        slide_related,
        slide_framework,
        slide_pipeline,
        slide_data,
        slide_corpus_table,
        slide_noise,
        slide_architecture,
        slide_hyperparams,
        slide_arabart,
        slide_metrics,
        slide_baselines,
        slide_training_results,
        slide_results_table,
        slide_cer_compare,
        slide_cer_dist,
        slide_confusion,
        slide_error_analysis,
        slide_discussion,
        slide_conclusion,
        slide_future,
        slide_refs,
        slide_thanks,
    ]
    for fn in builders:
        fn(prs)
        print("built:", fn.__name__)

    prs.save(str(OUT))
    print(f"\nSaved {OUT} with {len(prs.slides)} slides")


if __name__ == "__main__":
    build()
